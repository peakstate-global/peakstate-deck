#!/usr/bin/env python3
"""Export a deck built with this skill to PowerPoint, with EDITABLE native text.

    ~/.local/pptx-venv/bin/python export-pptx.py --deck /abs/path/to/index.html

--deck takes any deck that wires up the same ?export layout dump; the output is the deck folder's
name plus "-editable.pptx", written beside this script. Relative image sources resolve against the
DECK, so a deck anywhere on disk exports.

This lives with the slide rig rather than with any one deck, because every deck built on
deck-stage.js exports the same way and a copy per project goes stale.

The deck itself does the hard part: opened with ?export it measures its own laid-out geometry
and emits it as JSON (see dumpLayout in shift-keynote-deck.html). This script turns that JSON
into real PowerPoint shapes — rounded rectangles for the cards and bars, pictures for the two
SVG diagrams and the illustration, and text boxes you can edit in PowerPoint.

STATES AND HIDDEN SLIDES. A deck can carry states: several frames sharing one slide number,
where only one of them prints. PowerPoint has no such thing, so for now each state is exported
as its OWN slide, and the entry and exit animations are dropped — a state slide gets the Morph
transition instead, which is the closest PowerPoint has to the same idea. A slide the reviewer
hid is exported as a HIDDEN PowerPoint slide, so it is still in the file and still skipped in
the show. Presenter notes travel into the notes field either way.

Fidelity is close but not pixel-perfect: PowerPoint's line breaking differs from the browser's,
so a long heading may wrap a word earlier. For a pixel-perfect handout, print the HTML deck to
PDF instead (shift-keynote-deck.pdf) — this file exists so the text can be edited off this machine.
"""
import base64
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
# --deck <file> exports any deck that implements ?export/dumpLayout — the SHIFT deck, the W1-W13
# ways-of-working deck, and the consolidated great-AI-shifts deck all reuse the same measuring
# routine. Takes a filename or a relative path, so decks outside this folder work too.
DECK = HERE / (sys.argv[sys.argv.index("--deck") + 1] if "--deck" in sys.argv
               else "shift-keynote-deck.html")
if not DECK.exists():
    sys.exit(f"no deck at {DECK}. Pass --deck /abs/path/to/index.html")
# --public exports the public-sector cut of the SHIFT deck (?public): four inversions re-metered
# for an agency, three extra pressures appended. Other decks ignore it.
PUBLIC = "--public" in sys.argv
NAME = DECK.parent.name if DECK.stem == "index" else DECK.stem
OUT = HERE / (f"{NAME}-public-editable.pptx" if PUBLIC else f"{NAME}-editable.pptx")
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

# The deck's own coordinate space. Older decks (SHIFT, ways-of-working) are
# 1600x900 and uniformly dark, and their dumps say nothing about either, so these
# stay the fallback. A dump that carries "canvas" and per-slide "bg" — the
# consolidated Peak State deck does, at 1920x1080 across three surfaces —
# overrides them in main().
PX_W, PX_H = 1600, 900
EMU_PER_PX = int(Inches(13.333) / PX_W)  # 16:9 at 13.333in wide
PT_PER_PX = 72 * 13.333 / PX_W         # 1 deck px in points
BG = RGBColor(0x0D, 0x16, 0x26)        # --bg


def css_rgb(value: str, fallback: RGBColor) -> RGBColor:
    """'rgb(245, 240, 232)' / 'rgba(...)' / '#F5F0E8' -> RGBColor."""
    if not value:
        return fallback
    m = re.match(r"rgba?\(([^)]+)\)", value.strip())
    if m:
        parts = [p.strip() for p in m.group(1).replace("/", ",").split(",")]
        try:
            r, g, b = (int(round(float(p))) for p in parts[:3])
        except ValueError:
            return fallback
        # fully transparent means "whatever is behind", which for a slide is nothing
        if len(parts) > 3:
            try:
                if float(parts[3]) == 0:
                    return fallback
            except ValueError:
                pass
        return RGBColor(r, g, b)
    m = re.match(r"#([0-9a-fA-F]{6})$", value.strip())
    return RGBColor.from_string(m.group(1).upper()) if m else fallback
ALIGN = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}


SHAPE_HINTS = {
    "snip": MSO_SHAPE.SNIP_1_RECTANGLE,   # one cut corner — the products language
    "chevron": MSO_SHAPE.CHEVRON,         # notched left, point right — ways of working
    "pentagon": MSO_SHAPE.PENTAGON,       # flat left, point right — the head of a run
}


def flip_h(shape) -> None:
    """Mirror a shape horizontally. python-pptx exposes no API for this."""
    shape.element.spPr.xfrm.set("flipH", "1")


def strip_theme_style(shape) -> None:
    """Drop the <p:style> reference an AutoShape is born with.

    It carries <a:effectRef idx="2"> — the theme's drop shadow — which is
    reapplied even when the shape sets an empty <a:effectLst/>. Every card in
    this deck came out with a shadow the design system explicitly avoids. The
    fill, line and font are all set explicitly here anyway, so the whole
    reference is dead weight.
    """
    style = shape.element.find(qn("p:style"))
    if style is not None:
        shape.element.remove(style)


def pill(shape) -> None:
    """Push a rounded rectangle's corner radius to its maximum — a pill."""
    try:
        shape.adjustments[0] = 0.5
    except (IndexError, ValueError):
        pass


def rotate_text(shape) -> None:
    """Stand a text frame on end (bottom-to-top).

    A CSS writing-mode:vertical-rl box is tall and narrow; PowerPoint has no
    such thing, so the label came out as one letter per line. Rotating the body
    and swapping the box's own width and height about its centre puts it back.
    """
    tf = shape.text_frame
    tf._bodyPr.set("vert", "vert270")
    tf.word_wrap = False
    el = shape.element
    x, y = el.spPr.xfrm.off.x, el.spPr.xfrm.off.y
    cx, cy = el.spPr.xfrm.ext.cx, el.spPr.xfrm.ext.cy
    el.spPr.xfrm.off.x = int(x + (cx - cy) / 2)
    el.spPr.xfrm.off.y = int(y + (cy - cx) / 2)
    el.spPr.xfrm.ext.cx, el.spPr.xfrm.ext.cy = cy, cx


# ── The theme: layout names and font families ─────────────────────────────────
# A theme is optional. --theme <name|path|none>, or DECK_THEME in the
# environment. A bare name is looked up under <repo>/themes/<name>/theme.json.
# With no theme the exporter still runs: the layouts are called Dark and Light
# and no fonts are embedded.
GENERIC_FACES = {"serif", "sans-serif", "monospace", "cursive", "fantasy",
                 "system-ui", "ui-serif", "ui-sans-serif", "ui-monospace"}


def load_theme() -> dict:
    arg = (sys.argv[sys.argv.index("--theme") + 1] if "--theme" in sys.argv
           else os.environ.get("DECK_THEME", "peak-state"))
    if arg in ("", "none"):
        return {}
    given = Path(arg).expanduser()
    for cand in (given, given / "theme.json", HERE.parent / "themes" / arg / "theme.json"):
        if cand.is_file():
            return json.loads(cand.read_text())
    if "--theme" in sys.argv:
        sys.exit(f"no theme.json for --theme {arg}")
    return {}


THEME = load_theme()
_THEME_FONTS = THEME.get("fonts", {})

# The theme names its families; everything else in a CSS font stack is a
# fallback the browser only reaches when the real one is missing. Take the first
# name and keep it, so PowerPoint asks for the same face the deck drew with.
KNOWN_FACES = tuple(_THEME_FONTS.get("faces", ()))
FALLBACK_FACE = _THEME_FONTS.get("fallbackFace", "Helvetica Neue")


def face(css_font_family: str | None) -> str:
    """'Spectral, "Iowan Old Style", Georgia, serif' -> 'Spectral'."""
    if not css_font_family:
        return FALLBACK_FACE
    for part in css_font_family.split(","):
        name = part.strip().strip("\"'")
        for known in KNOWN_FACES:
            if name.lower() == known.lower():
                return known
    if not KNOWN_FACES:
        # No theme to match against, so trust the deck: the first real name in
        # the stack is the face it drew with.
        for part in css_font_family.split(","):
            name = part.strip().strip("\"'")
            if name and name.lower() not in GENERIC_FACES:
                return name
    return FALLBACK_FACE


# Naming the face is only half of it: a recipient without the font still sees a
# substitute. PowerPoint gives an embedded family four slots — regular, bold,
# italic, boldItalic — so the deck's finer weights (Spectral 200/300, Inter
# 500/600) are synthesised by the renderer from the nearest slot. Keep the fonts
# installed locally as well; embedding is for the people you send the file to.
# The fonts are referenced, never shipped in this repository. Three places are
# tried in order: DECK_FONT_DIR, a fonts/ folder at the repository root, and the
# older design-system path.
def font_dir() -> Path:
    candidates = [os.environ.get("DECK_FONT_DIR"),
                  HERE.parent / "fonts",
                  Path.home() / ".claude/skills/peak-state-design/fonts"]
    for c in candidates:
        if c and Path(c).expanduser().is_dir():
            return Path(c).expanduser()
    return HERE.parent / "fonts"


FONT_DIR = font_dir()
EMBED_FAMILIES = _THEME_FONTS.get("embed", {})
FONT_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"
SLOTS = (("regular", "Regular"), ("bold", "Bold"),
         ("italic", "Italic"), ("boldItalic", "BoldItalic"))


def embed_fonts(pptx_path: Path) -> None:
    """Rewrite the package in place with the three brand faces embedded."""
    import zipfile

    with zipfile.ZipFile(pptx_path) as zin:
        parts = {n: zin.read(n) for n in zin.namelist()}

    pres = parts["ppt/presentation.xml"].decode()
    if "embeddedFontLst" in pres:
        return

    rels = parts["ppt/_rels/presentation.xml.rels"].decode()
    next_id = max(int(n) for n in re.findall(r'Id="rId(\d+)"', rels)) + 1
    new_rels, font_list, n, missing = [], [], 0, []

    for family, pattern in EMBED_FAMILIES.items():
        entry = [f'<p:embeddedFont><p:font typeface="{family}" pitchFamily="18" charset="0"/>']
        for slot, suffix in SLOTS:
            src = FONT_DIR / pattern.format(suffix)
            if not src.exists():
                missing.append(src.name)
                continue
            n += 1
            parts[f"ppt/fonts/font{n}.fntdata"] = src.read_bytes()
            rid = f"rId{next_id}"
            next_id += 1
            new_rels.append(
                f'<Relationship Id="{rid}" Type="{FONT_REL}" Target="fonts/font{n}.fntdata"/>')
            entry.append(f'<p:{slot} r:id="{rid}"/>')
        entry.append("</p:embeddedFont>")
        if len(entry) > 2:
            font_list.append("".join(entry))

    if missing:
        print(f"  fonts not found, slots skipped: {', '.join(missing)}")
    if not font_list:
        print("  no fonts embedded")
        return

    parts["ppt/_rels/presentation.xml.rels"] = rels.replace(
        "</Relationships>", "".join(new_rels) + "</Relationships>").encode()

    ct = parts["[Content_Types].xml"].decode()
    if "fntdata" not in ct:
        # INSIDE <Types>. Matching the first ">" in the file lands on the end of
        # the <?xml?> declaration and puts the element outside the root.
        ct = re.sub(r"(<Types\b[^>]*>)",
                    r'\1<Default Extension="fntdata" ContentType="application/x-fontdata"/>',
                    ct, count=1)
        parts["[Content_Types].xml"] = ct.encode()

    # schema order: embeddedFontLst sits after notesSz, before defaultTextStyle
    pres = re.sub(r"(<p:notesSz[^>]*/>)",
                  r"\1<p:embeddedFontLst>" + "".join(font_list) + "</p:embeddedFontLst>",
                  pres, count=1)

    def set_attr(xml: str, name: str, value: str) -> str:
        """Set-or-replace. python-pptx's template already sets saveSubsetFonts,
        and a duplicated attribute makes the part unparseable."""
        if re.search(rf'<p:presentation\b[^>]*\b{name}="', xml):
            return re.sub(rf'({name}=")[^"]*(")', rf"\g<1>{value}\g<2>", xml, count=1)
        return xml.replace("<p:presentation ", f'<p:presentation {name}="{value}" ', 1)

    pres = set_attr(pres, "embedTrueTypeFonts", "1")
    pres = set_attr(pres, "saveSubsetFonts", "0")
    parts["ppt/presentation.xml"] = pres.encode()

    with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in parts.items():
            zout.writestr(name, data)
    print(f"  embedded {n} font files")


def accent_underline(run, hexcode: str) -> None:
    """A heavy underline in the accent colour, under text that keeps its own.

    The deck draws its one signal as a CSS background gradient, which has no
    computed property the dumper can read — so those five marks arrived in
    PowerPoint as run splits with no formatting at all, and the emphasis was
    silently gone. python-pptx exposes u= but not the underline fill, so the
    fill element is built by hand.
    """
    rPr = run.font._rPr
    rPr.set("u", "heavy")
    fill = rPr.makeelement(qn("a:uFill"), {})
    solid = rPr.makeelement(qn("a:solidFill"), {})
    clr = rPr.makeelement(qn("a:srgbClr"), {"val": hexcode})
    solid.append(clr)
    fill.append(solid)
    rPr.append(fill)


def bullet(para) -> None:
    """A real bullet glyph — python-pptx exposes no API for list formatting."""
    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", "171450")
    pPr.set("indent", "-171450")
    ch = pPr.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar", {"char": "\u2022"})
    pPr.append(ch)


def rgb(css, over=None):
    """CSS colour -> RGBColor.

    PowerPoint fills are opaque, so a translucent CSS fill must be composited
    against whatever sits behind it. Without `over`, rgba(169,61,26,.16) — a
    faint tint of sienna on sand — arrives as full-strength sienna and swallows
    the text sitting on it.
    """
    if not css:
        return None
    m = re.match(r"rgba?\(([\d.]+),\s*([\d.]+),\s*([\d.]+)(?:,\s*([\d.]+))?", css)
    if not m or (m.group(4) is not None and float(m.group(4)) == 0):
        return None
    r, g, b = (float(m.group(k)) for k in (1, 2, 3))
    a = float(m.group(4)) if m.group(4) is not None else 1.0
    if a < 1 and over is not None:
        r, g, b = (c * a + o * (1 - a) for c, o in zip((r, g, b), (over[0], over[1], over[2])))
    return RGBColor(*(int(round(v)) for v in (r, g, b)))


NS_MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"
NS_P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
NS_P159 = "http://schemas.microsoft.com/office/powerpoint/2015/main"
MORPH_MS = 700          # PowerPoint's morph is coarser than the deck's; 700 reads even

def morph(slide) -> None:
    """Give a slide the Morph transition, the closest PowerPoint has to a state change.

    python-pptx has no transition API, so this is the XML PowerPoint itself writes: the
    morph element lives in a 2015 namespace, wrapped in AlternateContent so an older
    reader falls back to a fade rather than refusing the file. Schema order puts the
    transition after clrMapOvr, which is where appending lands it.
    """
    P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    NS = {'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006',
          'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
          'p159': 'http://schemas.microsoft.com/office/powerpoint/2015/main'}
    xml = (
        '<mc:AlternateContent xmlns:mc="{mc}" xmlns:p="{p}">'
        '<mc:Choice xmlns:p159="{p159}" Requires="p159">'
        '<p:transition xmlns:p14="{p14}" spd="slow" p14:dur="{ms}">'
        '<p159:morph option="byObject"/>'
        '</p:transition>'
        '</mc:Choice>'
        '<mc:Fallback>'
        '<p:transition spd="slow"><p:fade/></p:transition>'
        '</mc:Fallback>'
        '</mc:AlternateContent>'
    ).format(ms=MORPH_MS, p=P, **NS)
    # PowerPoint writes these namespaces on the SLIDE ROOT and lists them as
    # ignorable. Without that, the reader is entitled to drop the whole
    # AlternateContent block, which is exactly what it was doing: the morph was
    # in the file, schema-valid, and never applied.
    root = slide._element
    for prefix, uri in (("p14", NS_P14), ("p159", NS_P159), ("mc", NS_MC)):
        root.set(f"{{http://www.w3.org/2000/xmlns/}}{prefix}", uri)
    IGNORABLE = f"{{{NS_MC}}}Ignorable"
    ignorable = root.get(IGNORABLE) or ""
    wanted = [t for t in ("p14", "p159") if t not in ignorable.split()]
    if wanted:
        root.set(IGNORABLE, " ".join(filter(None, [ignorable] + wanted)))
    root.append(parse_xml(xml))



# ── Layouts: the background and the branding live in the master ───────────────
# A deck built here runs on a small number of surfaces — this one has two, dark
# and light. Copying a full-bleed rectangle and a brand mark onto all 38 slides
# makes 38 things to keep in sync; putting them on a LAYOUT makes one. Change the
# dark colour in the master and every dark slide follows.
LAYOUT_NAMES = tuple(THEME.get("layoutNames", ("Dark", "Light")))
COLOUR_TOLERANCE = 8          # two surfaces this close are the same surface


def _near(a, b, tol=COLOUR_TOLERANCE):
    return all(abs(x - y) <= tol for x, y in zip(a, b))


def surfaces(slides):
    """The deck's own surfaces, commonest first, deduped within a tolerance."""
    counts = {}
    for s in slides:
        c = css_rgb(s.get("bg"), BG)
        key = (c[0], c[1], c[2])
        for seen in counts:
            if _near(seen, key):
                counts[seen] += 1
                break
        else:
            counts[key] = 1
    return [RGBColor(*k) for k, _ in sorted(counts.items(), key=lambda kv: -kv[1])]


def clear_layout(layout) -> None:
    """Strip a stock layout back to nothing, so what we add is all there is."""
    tree = layout.shapes._spTree
    for sp in list(tree):
        if sp.tag.endswith("}sp") or sp.tag.endswith("}pic") or sp.tag.endswith("}graphicFrame"):
            tree.remove(sp)


def adopt(layout, shape, source_slide, first=False):
    """Move a finished shape from a scratch slide onto a layout.

    python-pptx can read a layout's shapes but not add to them, so the shape is
    built where the library will build it and then moved. A picture also carries
    a relationship id that points into the SLIDE's part, so the image is related
    to the layout as well and the id rewritten, or the layout renders a grey box
    with a red cross.
    """
    el = shape._element
    el.getparent().remove(el)
    tree = layout.shapes._spTree
    tree.insert(2, el) if first else tree.append(el)      # 2 = after nvGrpSpPr/grpSpPr
    for blip in el.iter(qn("a:blip")):
        old_id = blip.get(qn("r:embed"))
        if not old_id:
            continue
        image_part = source_slide.part.related_part(old_id)
        blip.set(qn("r:embed"), layout.part.relate_to(image_part, RT.IMAGE))
    return el


def build_layouts(prs, colours):
    """One layout per surface, each with its own full-bleed background.

    python-pptx cannot mint a layout, so two stock ones are emptied and reused.
    Two is what this deck needs; a third surface keeps its own rectangle and is
    named in the summary rather than silently rendered on the wrong ground.
    """
    scratch = prs.slides.add_slide(prs.slide_layouts[6])
    out = {}
    # One layout per surface the deck actually uses, up to three. A surface with
    # no layout keeps its own rectangle and is named in the summary.
    stock = (6, 5, 4)
    for i, colour in enumerate(colours[:3]):
        lay = prs.slide_layouts[stock[i]]
        clear_layout(lay)
        # A theme need not name every surface a deck turns out to use.
        lay.name = LAYOUT_NAMES[i] if i < len(LAYOUT_NAMES) else f"Surface {i + 1}"
        bg = scratch.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                      prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = colour
        bg.line.fill.background()
        bg.shadow.inherit = False
        strip_theme_style(bg)
        adopt(lay, bg, scratch, first=True)
        out[(colour[0], colour[1], colour[2])] = lay
    return out, scratch


def lay_colour(layout):
    """The layout's own background colour, read back off the rectangle."""
    m = re.search(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', layout._element.xml)
    if not m:
        return (0, 0, 0)
    v = m.group(1)
    return (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def layout_for(layouts, colour):
    for key, layout in layouts.items():
        if _near(key, (colour[0], colour[1], colour[2])):
            return layout
    return None


def slide_number_field(shape) -> None:
    """Turn a text box into PowerPoint's own slide-number field.

    Typed text does not renumber when a slide is inserted; a field does. The
    deck marks which box holds the number, and this replaces its run with the
    field so the file behaves like one a person built.
    """
    from pptx.oxml.ns import qn as _qn
    r = shape.text_frame.paragraphs[0].runs[0]._r
    fld = parse_xml(
        '<a:fld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' id="{B7C7E4B4-9C8A-4D0B-9E3E-1F2A3B4C5D6E}" type="slidenum">'
        '<a:t>1</a:t></a:fld>')
    fld.insert(0, r.find(_qn("a:rPr")))
    parent = r.getparent()
    parent.insert(list(parent).index(r), fld)
    parent.remove(r)


def unmatched_morphs(slides):
    """Named shapes that have no partner on the neighbouring slide.

    PowerPoint pairs shapes for a morph by name. A name present on one slide of a
    pair and missing on the next is a morph that silently does not happen, which
    is exactly the class of defect that is invisible until it is on a screen in
    front of a room.
    """
    def names(s):
        return {i["morphName"] for i in s["items"] if i.get("morphName")}
    out = []
    for a, b in zip(slides, slides[1:]):
        if not (a.get("state") and a.get("state") == b.get("state")):
            continue
        for n in names(a) ^ names(b):
            out.append(f"{n} between slides {a['index']} and {b['index']}")
    return out


def name_for_morph(shape, item) -> None:
    """Give a shape the name the deck asked for.

    PowerPoint's morph pairs shapes between two slides by NAME, and a name that
    begins `!!` forces the pair even when the text inside differs. Without it a
    sentence exported as one shape has nothing to morph into, so it fades out and
    the next one fades in — which is what the deck was doing before the parts of
    the sentence were exported separately.
    """
    n = item.get("morphName")
    if n:
        shape._element._nvXxPr.cNvPr.set("name", n)


def hide(slide) -> None:
    """A hidden slide stays in the file and is skipped in the show."""
    slide._element.set('show', '0')


def layout() -> list[dict]:
    """The deck's own measurement of itself, taken in a real browser.

    Headless Chrome with --dump-dom was the original route and it became
    unreliable: it shares the running browser's profile, and its virtual-time
    budget expires on a page that has not finished booting, so the dump comes
    back either empty or not at all. Playwright drives its own browser and waits
    for the page to say it is ready, which is the difference between "usually"
    and "always".
    """
    script = r"""
    const { chromium } = require('playwright');
    (async () => {
      const [url, out] = process.argv.slice(2);
      const b = await chromium.launch();
      const p = await b.newPage({ viewport: { width: 1920, height: 1080 } });
      await p.goto(url, { waitUntil: 'networkidle' });
      // attached, not visible: the layout arrives as a <script> tag, which is
      // never visible and would time out on the default state.
      await p.waitForSelector('#layout', { state: 'attached', timeout: 60000 });
      const json = await p.$eval('#layout', el => el.textContent);
      require('fs').writeFileSync(out, json);
      await b.close();
    })();
    """
    # node resolves a require against the SCRIPT's folder, not the working
    # directory, so the runner is written where playwright is installed.
    node_root = Path(os.environ.get("DECK_NODE_ROOT", "/tmp"))
    assert (node_root / "node_modules" / "playwright").exists(), (
        f"playwright is not installed for node under {node_root}. Install it there, or set "
        f"DECK_NODE_ROOT to a folder that has it.")
    with tempfile.TemporaryDirectory(dir=node_root) as d:
        js, js_out = Path(d) / "dump.js", Path(d) / "layout.json"
        js.write_text(script, encoding="utf-8")
        url = DECK.as_uri() + ("?export&public" if PUBLIC else "?export")
        r = subprocess.run(["node", str(js), url, str(js_out)],
                           capture_output=True, text=True, cwd="/tmp")
        assert js_out.exists(), (
            "the deck did not emit its layout. Is ?export still wired up, and is "
            "playwright installed for node in /tmp?\n" + r.stderr[-600:])
        return json.loads(js_out.read_text(encoding="utf-8"))


def svg_styles() -> str:
    """Lift the deck's own SVG rules so a standalone render looks identical."""
    # Only the <style> block — the rest of the file carries a base64 image, and scanning it
    # with this regex takes minutes.
    doc = DECK.read_text(encoding="utf-8")
    css = re.search(r"<style>(.*?)</style>", doc, re.S).group(1)
    rules = re.findall(r"[^{}]*(?:\.chart|svg\.venn|svg\.swing)[^{}]*\{[^{}]*\}", css)
    body = "\n".join(rules)
    for var, hexcode in (("--m1", "#5b8ac9"), ("--m2", "#c9a227"), ("--m3", "#4aa6a0"),
                         ("--gold", "#d1b561"), ("--fg", "#eceadf"), ("--muted", "#93a0b5"),
                         ("--line", "#223049"), ("--danger", "#e07b74")):
        body = body.replace(f"var({var})", hexcode)
    # font-family: inherit in the deck's own rules resolves to a serif default in a standalone
    # file, so the override has to win on specificity.
    return ("<style>" + body +
            "\ntext, tspan { font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif !important; }"
            "</style>")


def symbol_defs(markup: str) -> str:
    """The <symbol> a <use href="#id"> points at, lifted out of the deck.

    An inline SVG that references a symbol defined elsewhere in the page renders
    as nothing once it is written to a file of its own — which is why the brand
    diamond was missing from the export while being present on every slide.
    """
    ids = re.findall(r'href="#([\w-]+)"', markup) + re.findall(r'xlink:href="#([\w-]+)"', markup)
    if not ids:
        return ""
    page = DECK.read_text(encoding="utf-8")
    out = []
    for i in ids:
        m = re.search(rf'<symbol[^>]*id="{re.escape(i)}".*?</symbol>', page, re.S)
        if m:
            out.append(m.group(0))
    return "<defs>" + "".join(out) + "</defs>" if out else ""


def raster(item: dict, tmp: Path, name: str, recolour: RGBColor | None = None) -> Path:
    """SVG diagrams and the illustration ride along as pictures — only the text is editable."""
    if item["type"] == "img":
        src = item["src"]
        head, _, b64 = src.partition(",")
        out = tmp / f"{name}.png"
        # A relative src is relative to the DECK, not to this script. They are the
        # same folder for the decks that live here and different for any other.
        out.write_bytes(base64.b64decode(b64) if "base64" in head
                        else (DECK.parent / src).read_bytes())
        return out
    svg = tmp / f"{name}.svg"
    # A standalone SVG has none of the deck's stylesheet, so its class-based rules (stroke
    # widths, fills, fonts) must travel with it — otherwise the curves render as filled blobs.
    markup = item["svg"]
    if "xmlns" not in markup:  # inline SVG inherits the namespace from HTML; a file cannot
        markup = markup.replace("<svg ", '<svg xmlns="http://www.w3.org/2000/svg" ', 1)
    markup = markup.replace(">", ">" + svg_styles() + symbol_defs(markup), 1)
    if recolour is not None:
        # One ground needs one ink. The mark is drawn in the deck's light colour,
        # which vanishes on a light layout.
        hexcode = "#%02X%02X%02X" % (recolour[0], recolour[1], recolour[2])
        markup = markup.replace("</svg>",
                                f'<style>*{{fill:{hexcode}!important;'
                                f'stroke:{hexcode}!important}}</style></svg>', 1)
    svg.write_text(
        '<?xml version="1.0"?>' + markup
        .replace("var(--m1)", "#5b8ac9").replace("var(--m2)", "#c9a227")
        .replace("var(--m3)", "#4aa6a0").replace("var(--gold)", "#d1b561")
        .replace("var(--fg)", "#eceadf").replace("var(--muted)", "#93a0b5")
        .replace("var(--line)", "#223049"), encoding="utf-8")
    png = tmp / f"{name}.png"
    subprocess.run([CHROME, "--headless=new", "--disable-gpu", "--default-background-color=00000000",
                    f"--window-size={int(item['w'])},{int(item['h'])}",
                    "--force-device-scale-factor=2", "--virtual-time-budget=3000",
                    f"--screenshot={png}", svg.as_uri()], check=True, capture_output=True)
    return png


def main() -> None:
    global PX_W, PX_H, EMU_PER_PX, PT_PER_PX
    slides = layout()
    canvas = (slides[0] or {}).get("canvas") if slides else None
    if canvas:
        PX_W, PX_H = canvas["w"], canvas["h"]
        EMU_PER_PX = int(Inches(13.333) / PX_W)
        PT_PER_PX = 72 * 13.333 / PX_W
    print(f"{len(slides)} slides at {PX_W}x{PX_H}")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    found = surfaces(slides)
    layouts, scratch = build_layouts(prs, found)
    show_brand = bool(slides and slides[0].get("branding", True))
    _grounds = [css_rgb(x.get("bg"), BG) for x in slides]
    darkest_ink, lightest_ink = min(_grounds, key=sum), max(_grounds, key=sum)
    report = {"layouts": {}, "own_background": [], "brand_on_layout": 0,
              "surfaces": len(found), "numbers": 0}
    brand_done = False
    brand_items = []
    number_boxes = []
    tmp = Path(tempfile.mkdtemp(prefix="deck-native-"))

    try:
        for s in slides:
            slide_bg = css_rgb(s.get("bg"), BG)
            # `layout` is already the name of the function that dumps the deck's
            # geometry; this is the PowerPoint one.
            ground = layout_for(layouts, slide_bg)
            slide = prs.slides.add_slide(ground or blank)
            if ground is None:
                # A third surface. It keeps its own rectangle and gets named at
                # the end, rather than being drawn on the wrong ground quietly.
                bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                            prs.slide_width, prs.slide_height)
                bg.fill.solid()
                bg.fill.fore_color.rgb = slide_bg
                bg.shadow.inherit = False
                bg.line.fill.background()
                strip_theme_style(bg)
                report["own_background"].append(s["index"])
            else:
                report["layouts"][ground.name] = report["layouts"].get(ground.name, 0) + 1

            for n, item in enumerate(s["items"]):
                brand = item.get("brand")
                if brand == "mark":
                    # Held back and drawn once onto each layout, below. Identical
                    # on every ordinary slide, so a copy per slide is 38 things
                    # to sync. The title slide's own mark is marked "title" and
                    # is NOT this: it sits top right and belongs to that slide.
                    if show_brand and not brand_done:
                        # Keep gathering until the mark is COMPLETE. The title
                        # slide shows the diamond and hides the words, so a
                        # collector that stops at the first slide carrying any
                        # mark takes the diamond alone and the wordmark never
                        # reaches a layout.
                        if not any(b["type"] == item["type"] for b in brand_items):
                            brand_items.append(item)
                    continue
                if brand == "title":
                    pass          # drawn on its own slide, like any other item
                if brand == "number" and not show_brand:
                    continue
                x, y = Emu(int(item["x"] * EMU_PER_PX)), Emu(int(item["y"] * EMU_PER_PX))
                w, h = Emu(int(item["w"] * EMU_PER_PX)), Emu(int(item["h"] * EMU_PER_PX))

                if item["type"] == "rule":
                    thick = Emu(int(max(1, item.get("borderW", 1)) * EMU_PER_PX))
                    edge = item.get("edge")
                    horizontal = edge in ("top", "bottom")
                    # A bottom border draws at the FOOT of its box, not its head.
                    top = Emu(int(y) + int(h) - int(thick)) if edge == "bottom" else y
                    line = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, x, top,
                        w if horizontal else thick, thick if horizontal else h)
                    line.fill.solid()
                    line.fill.fore_color.rgb = (rgb(item.get("border"), slide_bg)
                                               or RGBColor(0x22, 0x30, 0x49))
                    line.line.fill.background()
                    line.shadow.inherit = False
                    strip_theme_style(line)
                    name_for_morph(line, item)

                elif item["type"] == "box":
                    # A deck built on clip-path loses its whole shape language in
                    # PowerPoint unless the silhouettes are mapped to AutoShapes.
                    # The deck names the one it wants via data-pptx-shape.
                    hint = item.get("shape")
                    if item.get("oval"):
                        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
                    elif item.get("pill"):
                        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
                        pill(shape)
                    elif hint:
                        shape = slide.shapes.add_shape(SHAPE_HINTS[hint], x, y, w, h)
                        # CHEVRON and PENTAGON size their notch as a fraction of
                        # the shape, so segments of different heights ended up
                        # with different notch depths and the run stopped lining
                        # up. The deck states the depth it drew; hold it.
                        notch = item.get("notch") or 0
                        if notch and item.get("h"):
                            try:
                                shape.adjustments[0] = min(0.5, notch / item["h"])
                            except (IndexError, ValueError):
                                pass
                        if hint == "snip":
                            # SNIP_1_RECTANGLE snips the top-RIGHT corner; the deck
                            # snips the top-left. Mirror it.
                            flip_h(shape)
                    else:
                        shape = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE if item.get("radius", 0) > 2
                            else MSO_SHAPE.RECTANGLE, x, y, w, h)
                    fill = rgb(item.get("fill"), slide_bg)
                    if fill:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = fill
                    else:
                        shape.fill.background()
                    border = rgb(item.get("border"), slide_bg)
                    if border:
                        shape.line.color.rgb = border
                        shape.line.width = Pt(max(0.75, item.get("borderW", 1) * PT_PER_PX))
                    else:
                        shape.line.fill.background()
                    shape.shadow.inherit = False
                    strip_theme_style(shape)

                elif item["type"] in ("svg", "img"):
                    path = raster(item, tmp, f"s{s['index']}-{n}")
                    pic = slide.shapes.add_picture(str(path), x, y, width=w, height=h)
                    # The browser uses object-fit: cover, so crop the overflow rather than
                    # stretching the source into a different aspect ratio.
                    src = pic.image.size[0] / pic.image.size[1]
                    box = item["w"] / item["h"]
                    if src > box:      # source is wider — trim the sides
                        trim = (1 - box / src) / 2
                        pic.crop_left = pic.crop_right = trim
                    elif box > src:    # source is taller — trim top and bottom
                        trim = (1 - src / box) / 2
                        pic.crop_top = pic.crop_bottom = trim

                else:
                    # PowerPoint measures type wider than the browser does, so a
                    # box cut exactly to the browser's width clips or wraps a
                    # short string. A split shape is always a fragment — a word, a
                    # percentage — so it gets slack rather than a wrap.
                    if item.get("morphName"):
                        # A split shape is a FRAGMENT of a sentence, never a
                        # paragraph. PowerPoint measures type wider than the
                        # browser, so a box cut to the browser's width wraps a
                        # fragment in half. Slack, and never wrap.
                        w = Emu(int(int(w) * 1.3))
                    box = slide.shapes.add_textbox(x, y, w, h)
                    name_for_morph(box, item)
                    if item.get("brand") == "number":
                        number_boxes.append(box)
                    tf = box.text_frame
                    # PowerPoint substitutes fonts and measures wider than the
                    # browser, so a short, deliberately-unbreakable string — a
                    # big stat like "78%" — wraps mid-number and collides with
                    # its own caption. Nothing that short was ever meant to wrap.
                    tf.word_wrap = (not item.get("morphName")
                                    and len(item.get("text", "").split()) > 2)
                    pad = item.get("pad") or [0, 0, 0, 0]
                    tf.margin_top, tf.margin_right, tf.margin_bottom, tf.margin_left = (
                        Emu(int((v or 0) * EMU_PER_PX)) for v in pad)
                    # Middle-anchored: identical to top for a paragraph whose box is its own
                    # height, and it centres the short labels that sit inside a padded pill.
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    if item.get("rotate") == 270:
                        rotate_text(box)
                    pt = item["size"] * PT_PER_PX
                    if pt >= 30:  # PowerPoint breaks lines later than the browser does
                        pt *= 0.94
                    # Runs carry inline emphasis (the gold <em class="hl">); fall back to one
                    # run of the block's own style when the deck didn't record them.
                    runs = item.get("runs") or [{"text": item["text"],
                                                 "colour": item.get("colour"),
                                                 "weight": item.get("weight"),
                                                 "italic": item.get("italic"),
                                                 "strike": item.get("strike"),
                                                 "font": item.get("font"),
                                                 "underline": item.get("underline")}]
                    para = tf.paragraphs[0]
                    para.alignment = ALIGN.get(item.get("align"), PP_ALIGN.LEFT)
                    if item.get("bullet"):
                        bullet(para)
                    first = True
                    for spec in runs:
                        for k, chunk in enumerate(spec["text"].split("\n")):
                            if k:
                                para = tf.add_paragraph()
                                para.alignment = ALIGN.get(item.get("align"), PP_ALIGN.LEFT)
                            if not chunk:
                                continue
                            run = para.add_run()
                            run.text = chunk.upper() if item.get("caps") else chunk
                            if first and not run.text.startswith(" "):
                                first = False
                            run.font.size = Pt(round(pt, 1))
                            run.font.bold = (spec.get("weight") or 400) >= 600
                            run.font.italic = bool(spec.get("italic"))
                            run.font.name = face(spec.get("font") or item.get("font"))
                            if spec.get("underline"):
                                accent_underline(run, spec["underline"])
                            if spec.get("strike"):  # python-pptx has no strike property
                                run.font._rPr.set("strike", "sngStrike")
                            colour = rgb(spec.get("colour")) or rgb(item.get("colour"))
                            if colour:
                                run.font.color.rgb = colour

            # Done once the mark has both halves: the diamond and the words.
            if len({b["type"] for b in brand_items}) >= 2:
                brand_done = True
            slide.notes_slide.notes_text_frame.text = s["note"]
            if s.get("state"):
                morph(slide)
            if s.get("hidden"):
                hide(slide)
            print(".", end="", flush=True)

        # The branding, once per layout. A layout is a slide's ground, so what
        # sits here appears on every slide built on it and on none of them twice.
        for item in brand_items:
            for lay in layouts.values():
                x = Emu(int(item["x"] * EMU_PER_PX)); y = Emu(int(item["y"] * EMU_PER_PX))
                w = Emu(int(item["w"] * EMU_PER_PX)); h = Emu(int(item["h"] * EMU_PER_PX))
                # The mark is lifted from ONE slide, so whatever ink it wore
                # there is wrong for at least one layout. Ink is chosen per
                # ground instead of inherited: the deck's lightest surface on a
                # dark ground, its darkest on a light one.
                dark_ground = sum(lay_colour(lay)) < 384
                ink = lightest_ink if dark_ground else darkest_ink
                if item["type"] in ("img", "svg"):
                    png = raster(item, tmp, f"brand-{lay.name}-{item['type']}", recolour=ink)
                    pic = scratch.shapes.add_picture(str(png), x, y, w, h)
                    adopt(lay, pic, scratch)
                    report["brand_on_layout"] += 1
                    continue
                else:
                    box = scratch.shapes.add_textbox(x, y, w, h)
                    tf = box.text_frame
                    tf.word_wrap = False
                    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                    para = tf.paragraphs[0]
                    run = para.add_run()
                    run.text = item.get("text", "")
                    run.font.size = Pt(item["size"] * PT_PER_PX)
                    run.font.name = face(item.get("font"))
                    run.font.color.rgb = ink
                    adopt(lay, box, scratch)
                report["brand_on_layout"] += 1

        # The page number becomes PowerPoint's own field, so inserting a slide
        # renumbers the rest instead of leaving typed text behind.
        for box in number_boxes:
            try:
                slide_number_field(box)
                report["numbers"] += 1
            except Exception:
                pass
        # The scratch slide existed only so the library had somewhere to build
        # shapes that now live on the layouts.
        xml_slides = prs.slides._sldIdLst
        for sldId in list(xml_slides):
            if prs.slides.get(sldId.rId) is None or sldId.rId == scratch.part.partname:
                pass
        drop = [sldId for sldId in xml_slides
                if prs.part.related_part(sldId.rId) is scratch.part]
        for sldId in drop:
            prs.part.drop_rel(sldId.rId)
            xml_slides.remove(sldId)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

    prs.save(OUT)
    # Ticket 04: say what the run did, so a wrong assumption is visible here
    # rather than in the room.
    print()
    print(f"  surfaces found      {report['surfaces']}")
    for name, n in sorted(report["layouts"].items()):
        print(f"  on layout {name:<12} {n} slide(s)")
    if report["own_background"]:
        print(f"  kept their own background: slides {report['own_background']} "
              f"(their surface matched no layout)")
    print(f"  branding            {'on' if show_brand else 'OFF'}, "
          f"{report['brand_on_layout']} shape(s) placed on the layouts")
    print(f"  page numbers        {report['numbers']} field(s), they renumber themselves")
    print(f"  hidden slides       {sum(1 for s in slides if s.get('hidden'))}")
    print(f"  morph transitions   {sum(1 for s in slides if s.get('state'))}")
    print(f"  notes carried       {sum(1 for s in slides if (s.get('note') or '').strip())}")
    lonely = sorted(unmatched_morphs(slides))  # sorted: the finder returns a set
    if lonely:
        print(f"  WARNING: a named morph shape appears on one slide of a pair and not the "
              f"other, so that morph will not happen: {lonely}")
    embed_fonts(OUT)
    print(f"\n{OUT}")


if __name__ == "__main__":
    sys.exit(main())
